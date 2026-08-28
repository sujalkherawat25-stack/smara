"""Authenticated, bounded attachment storage for hosted Smara.

Attachments are transient conversation context, not long-term memory.  Files
are namespaced by account, given opaque ids, and stored below the mounted
Smara data volume.  The API enforces the 100 MB per-file and 150 MB per-upload
batch limits before a chat can reference an id.
"""
from __future__ import annotations

import hashlib
import io
import json
import mimetypes
import base64
import re
import secrets
from datetime import datetime, timezone
from pathlib import Path
from typing import Any

MAX_FILE_BYTES = 100 * 1024 * 1024
MAX_BATCH_BYTES = 150 * 1024 * 1024
MAX_ATTACHMENTS_PER_BATCH = 10
MAX_CONTEXT_CHARS = 120_000
# Keep multimodal requests bounded even when the upload itself is allowed to
# be large.  The original file remains available for future local tooling.
MAX_IMAGE_INLINE_BYTES = 12 * 1024 * 1024

_SAFE_NAME = re.compile(r"[^A-Za-z0-9._ -]+")
_TEXT_EXTENSIONS = {
    ".txt", ".md", ".markdown", ".csv", ".json", ".xml", ".html", ".htm",
    ".yaml", ".yml", ".log", ".py", ".js", ".ts", ".tsx", ".jsx", ".css",
    ".sql", ".toml", ".ini", ".env",
}


def _extract_text(path: str, filename: str, content_type: str) -> str:
    """Best-effort extraction for common office/document formats.

    Binary files remain safely attached and are reported as metadata; a
    missing optional parser never turns a successful upload into a 500.
    """
    data = Path(path).read_bytes()
    ext = Path(filename).suffix.lower()
    if content_type.startswith("text/") or ext in _TEXT_EXTENSIONS:
        return data.decode("utf-8", errors="ignore")
    try:
        if ext == ".pdf":
            from pypdf import PdfReader
            reader = PdfReader(io.BytesIO(data))
            return "\n\n".join((page.extract_text() or "").strip() for page in reader.pages)
        if ext == ".docx":
            import docx
            document = docx.Document(io.BytesIO(data))
            return "\n".join(p.text.strip() for p in document.paragraphs if p.text.strip())
        if ext == ".xlsx":
            import openpyxl
            workbook = openpyxl.load_workbook(io.BytesIO(data), read_only=True, data_only=True)
            try:
                rows: list[str] = []
                for sheet in workbook.worksheets:
                    rows.append(f"## Sheet: {sheet.title}")
                    for row in sheet.iter_rows(values_only=True):
                        values = ["" if value is None else str(value) for value in row]
                        if any(values): rows.append(" | ".join(values))
                return "\n".join(rows)
            finally:
                workbook.close()
        if ext == ".pptx":
            from pptx import Presentation
            presentation = Presentation(io.BytesIO(data))
            rows = []
            for number, slide in enumerate(presentation.slides, 1):
                rows.append(f"## Slide {number}")
                for shape in slide.shapes:
                    if getattr(shape, "has_text_frame", False) and shape.text.strip():
                        rows.append(shape.text.strip())
            return "\n".join(rows)
    except Exception:
        return ""
    return ""


def safe_filename(name: str | None) -> str:
    value = Path(name or "attachment").name.strip() or "attachment"
    value = _SAFE_NAME.sub("_", value)
    return value[:240]


def _now() -> str:
    return datetime.now(timezone.utc).isoformat()


class AttachmentStore:
    def __init__(self, root: str | Path):
        self.root = Path(root)
        self.root.mkdir(parents=True, exist_ok=True)

    def _account_dir(self, account_id: str) -> Path:
        # account ids are generated server-side; keep a defensive filename
        # boundary even if a future auth adapter changes their shape.
        directory = self.root / re.sub(r"[^A-Za-z0-9_-]", "_", account_id)
        directory.mkdir(parents=True, exist_ok=True)
        return directory

    async def save(self, account_id: str, upload: Any, *, size_limit: int = MAX_FILE_BYTES) -> dict[str, Any]:
        filename = safe_filename(getattr(upload, "filename", None))
        content_type = str(getattr(upload, "content_type", None) or mimetypes.guess_type(filename)[0] or "application/octet-stream")[:160]
        attachment_id = f"att_{secrets.token_urlsafe(18).replace('-', '_').replace('.', '_')}"
        directory = self._account_dir(account_id)
        data_path = directory / f"{attachment_id}.bin"
        metadata_path = directory / f"{attachment_id}.json"
        digest = hashlib.sha256()
        size = 0
        try:
            with data_path.open("wb") as target:
                while True:
                    chunk = await upload.read(1024 * 1024)
                    if not chunk:
                        break
                    size += len(chunk)
                    if size > size_limit:
                        raise ValueError(f'"{filename}" is over the 100 MB per-file limit.')
                    digest.update(chunk)
                    target.write(chunk)
            if size == 0:
                raise ValueError(f'"{filename}" is empty.')
            record = {
                "id": attachment_id,
                "account_id": account_id,
                "filename": filename,
                "content_type": content_type,
                "size": size,
                "sha256": digest.hexdigest(),
                "created_at": _now(),
            }
            metadata_path.write_text(json.dumps(record), encoding="utf-8")
            return record
        except Exception:
            data_path.unlink(missing_ok=True)
            metadata_path.unlink(missing_ok=True)
            raise

    def get(self, account_id: str, attachment_id: str) -> dict[str, Any] | None:
        if not re.fullmatch(r"att_[A-Za-z0-9_]+", attachment_id or ""):
            return None
        path = self._account_dir(account_id) / f"{attachment_id}.json"
        try:
            record = json.loads(path.read_text(encoding="utf-8"))
        except (OSError, json.JSONDecodeError):
            return None
        if record.get("account_id") != account_id:
            return None
        data_path = path.with_suffix(".bin")
        if not data_path.is_file():
            return None
        record["path"] = str(data_path)
        return record

    def delete(self, account_id: str, attachment_id: str) -> None:
        if not re.fullmatch(r"att_[A-Za-z0-9_]+", attachment_id or ""):
            return
        directory = self._account_dir(account_id)
        (directory / f"{attachment_id}.json").unlink(missing_ok=True)
        (directory / f"{attachment_id}.bin").unlink(missing_ok=True)

    def context_for(self, account_id: str, attachment_ids: list[str]) -> tuple[str, list[dict[str, Any]]]:
        """Return bounded text context and metadata for a chat turn."""
        parts: list[str] = []
        records: list[dict[str, Any]] = []
        remaining = MAX_CONTEXT_CHARS
        for attachment_id in attachment_ids[:MAX_ATTACHMENTS_PER_BATCH]:
            record = self.get(account_id, attachment_id)
            if not record:
                continue
            records.append({k: record[k] for k in ("id", "filename", "content_type", "size", "sha256")})
            filename = record["filename"]
            text = ""
            try:
                text = _extract_text(record["path"], filename, record["content_type"])[:remaining]
            except OSError:
                text = ""
            if text:
                parts.append(f"Attachment: {filename}\n{text}")
                remaining -= len(text)
            else:
                parts.append(
                    f"Attachment: {filename} ({record['content_type']}, {record['size']} bytes). "
                    "The file is attached but has no text preview in this turn."
                )
            if remaining <= 0:
                break
        return "\n\n".join(parts)[:MAX_CONTEXT_CHARS], records

    def image_inputs(self, account_id: str, attachment_ids: list[str]) -> list[dict[str, str]]:
        """Return small uploaded images as OpenAI-compatible data URLs.

        This is deliberately opt-in at the runtime: only a profile declared
        with ``capability=vision`` receives these bytes.  Large images stay in
        account storage and are described by the normal attachment context,
        avoiding accidental oversized provider requests.
        """
        images: list[dict[str, str]] = []
        for attachment_id in attachment_ids[:MAX_ATTACHMENTS_PER_BATCH]:
            record = self.get(account_id, attachment_id)
            if not record or int(record.get("size", 0)) > MAX_IMAGE_INLINE_BYTES:
                continue
            content_type = str(record.get("content_type") or "")
            if not content_type.startswith("image/"):
                continue
            try:
                encoded = base64.b64encode(Path(record["path"]).read_bytes()).decode("ascii")
            except (OSError, KeyError):
                continue
            images.append({
                "filename": str(record.get("filename", "image")),
                "data_url": f"data:{content_type};base64,{encoded}",
            })
        return images
